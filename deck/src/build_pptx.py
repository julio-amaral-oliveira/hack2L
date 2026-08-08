"""Gera BrandLoop-deck.pptx a partir do mesmo conteúdo do deck HTML.
Formato 16:9 (13,333 x 7,5 pol). Importa direto no Google Slides.
"""
import sys

from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PRETO   = C(0x0A, 0x0A, 0x0A)
CRU     = C(0xF3, 0xEF, 0xE5)
AMARELO = C(0xFF, 0xD4, 0x00)
CORAL   = C(0xFF, 0x66, 0x4A)
MENTA   = C(0x9F, 0xF5, 0xD0)
VIOLETA = C(0xA7, 0x8B, 0xFA)
CINZA   = C(0x8A, 0x8A, 0x86)
DIM     = C(0xA8, 0xA4, 0x9B)
CARD    = C(0x16, 0x16, 0x16)

SANS = "Arial"
MONO = "Courier New"

W, H = 13.333, 7.5
M = 0.75  # margem

LOGO = sys.argv[1]  # PNG da logo da Gorilla
OUT = sys.argv[2]   # caminho do .pptx de saída

prs = Presentation()
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[6]


def slide(bg=PRETO):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, In(W), In(H))
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def text(s, txt, x, y, w, h, size=14, color=CRU, bold=False, font=SANS,
         space=0, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=None):
    """Cada item de `txt` (str ou lista de (texto, cor, bold)) vira um parágrafo."""
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    itens = txt if isinstance(txt, list) else [txt]
    for i, item in enumerate(itens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        runs = item if isinstance(item, list) else [(item, color, bold)]
        for t, c, b in runs:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = b
            r.font.color.rgb = c; r.font.name = font
            if space:
                # espaçamento entre letras, para rótulos em mono
                from pptx.oxml.ns import qn
                r.font._rPr.set('spc', str(int(space * 100)))
    return tb


def card(s, x, y, w, h, fill=CARD):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(x), In(y), In(w), In(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.fill.background(); r.shadow.inherit = False
    return r


def label(s, txt, x, y, w, color=CINZA, size=9):
    return text(s, txt.upper(), x, y, w, 0.25, size=size, color=color, font=MONO, space=1.2)


# ---------------------------------------------------------------- 1 · tensão
s = slide()
label(s, "Hack2L · AI Agents · Canastra Ventures", M, 0.6, 8)
text(s, [[("A IA ensinou todo mundo a fazer anúncio.", CRU, True)],
         [("E agora todo anúncio parece igual.", AMARELO, True)]],
     M, 2.6, W - 2 * M, 2.6, size=48, line=1.0)
s.notes_slide.notes_text_frame.text = "0:00–0:25 — Abrir pela tensão. Não explicar demais."

# ------------------------------------------------- 2 · o sistema de agentes
s = slide()
label(s, "O que a BrandLoop é", M, 0.55, 8)
text(s, [[("Não é um prompt.", CRU, True)],
         [("São seis agentes com contrato entre eles.", CRU, True)]],
     M, 1.0, W - 2 * M, 1.3, size=28, line=1.05)
AG = [
    ("01", "Contexto", "Lê os arquivos da empresa e extrai fatos sob schema — não pode inventar o que não está no texto.", "/api/ingest"),
    ("02", "Entrevista", "Decide sozinho a próxima pergunta olhando quais campos do diagnóstico ainda estão vazios.", "/api/interview"),
    ("03", "Mercado", "Busca conversa real e devolve o diagnóstico preenchido, com a URL de cada evidência.", "/api/research"),
    ("04", "Roteiro", "Escreve sob a rubrica de Schwartz como restrição tipada — sem prova, ele reformula.", "/api/copy"),
    ("05", "Criativo", "Gera o vídeo vertical e o áudio com a voz da marca, por provider plugável.", "/api/video"),
    ("06", "Anúncio e funil", "Publica, lê a performance e decide onde mexer, dentro de um envelope de verba.", "/api/publish"),
]
cw, ch, gap = 3.85, 1.85, 0.22
for i, (n, ttl, desc, rota) in enumerate(AG):
    x = M + (i % 3) * (cw + gap)
    y = 2.65 + (i // 3) * (ch + gap)
    card(s, x, y, cw, ch)
    text(s, n, x + 0.22, y + 0.18, 1, 0.2, size=9, color=AMARELO, font=MONO, space=1.2)
    text(s, ttl, x + 0.22, y + 0.42, cw - 0.44, 0.3, size=14, bold=True)
    text(s, desc, x + 0.22, y + 0.78, cw - 0.44, 0.7, size=9.5, color=DIM, line=1.25)
    text(s, rota, x + 0.22, y + ch - 0.34, cw - 0.44, 0.2, size=8, color=CINZA, font=MONO)
text(s, "Nenhuma etapa espera humano para começar a próxima.",
     M, 6.7, W - 2 * M, 0.5, size=13, color=AMARELO, line=1.3)
s.notes_slide.notes_text_frame.text = "0:25–0:50 — É o slide que responde ao critério do hackathon. Onde a IA trabalha."

# ------------------------------------------------------- 3 · começa na marca
s = slide()
label(s, "Primeiro passo — a inversão", M, 0.6, 8)
text(s, "A BrandLoop começa onde as outras IAs terminam: na marca.",
     M, 1.05, 10.5, 1.4, size=30, bold=True, line=1.05)
text(s, "Antes de escrever uma linha, o agente entrevista as lideranças da empresa. Uma pergunta "
        "por vez, cada uma mirando só o campo que ainda falta — porque o material institucional diz "
        "o que a empresa gostaria de ser, e quem vende sabe qual objeção mata o negócio na terça-feira.",
     M, 2.7, 10.6, 1.2, size=13, color=DIM, line=1.35)
DIMS = [("Tom", "Como a marca fala quando é ela mesma", CORAL),
        ("Público", "Quem compra, e qual objeção trava a venda", MENTA),
        ("Estética", "O que a faz reconhecível sem o logo", VIOLETA)]
cw2 = 3.85
for i, (k, v, cor) in enumerate(DIMS):
    x = M + i * (cw2 + 0.22)
    card(s, x, 4.35, cw2, 1.7)
    text(s, k.upper(), x + 0.25, 4.6, cw2 - 0.5, 0.3, size=11, color=cor, bold=True, font=MONO, space=1.5)
    text(s, v, x + 0.25, 5.05, cw2 - 0.5, 0.8, size=12, line=1.3)
s.notes_slide.notes_text_frame.text = "0:50–1:15 — A entrevista com as lideranças vem ANTES da pesquisa de mercado."

# ----------------------------------------------------- 4 · vai ao mercado
s = slide()
label(s, "Só depois — a verificação", M, 0.6, 8)
text(s, "Aí sim ele confere lá fora se é isso mesmo.",
     M, 1.05, 10.5, 1.0, size=30, bold=True, line=1.05)
card(s, M, 2.35, W - 2 * M, 1.35)
try:
    s.shapes.add_picture(LOGO, In(M + 0.3), In(2.6), height=In(0.85))
except Exception:
    pass
text(s, "ANÁLISE DE MERCADO", M + 1.35, 2.62, 5, 0.25, size=9, color=CINZA, font=MONO, space=1.5)
text(s, "Reddit · X · Bluesky · LinkedIn · YouTube — conversa real, em tempo real",
     M + 1.35, 2.95, 9.5, 0.5, size=13, color=DIM)
MET = [("1.516", "conversas reais lidas"), ("5", "plataformas"), ("< 3 min", "menos que este pitch")]
cw3 = 3.85
for i, (v, k) in enumerate(MET):
    x = M + i * (cw3 + 0.22)
    card(s, x, 3.95, cw3, 1.5)
    text(s, v, x + 0.25, 4.2, cw3 - 0.5, 0.7, size=34, bold=True)
    text(s, k.upper(), x + 0.25, 4.95, cw3 - 0.5, 0.3, size=9, color=CINZA, font=MONO, space=1.2)
text(s, [[("E o que ele concluiu: ", CRU, True),
          ("esse mercado não quer mais uma promessa de venda — quer parar de fazer a conta de madrugada.", DIM, False)]],
     M, 5.85, W - 2 * M, 0.6, size=13, line=1.35)
s.notes_slide.notes_text_frame.text = "1:15–1:40 — O mérito é a CONCLUSÃO, não a infraestrutura de busca."

# --------------------------------------------------- 5 · criativo + funil
s = slide()
label(s, "Identidade vira criativo — e o criativo persegue a venda", M, 0.55, 10)
cwm = (W - 2 * M - 0.22) / 2
card(s, M, 1.0, cwm, 1.75)
text(s, "QUALQUER IA · MESMO MODELO", M + 0.25, 1.2, cwm - 0.5, 0.25, size=9, color=CINZA, font=MONO, space=1.2)
text(s, "“Cadastre seu restaurante no iFood e venda mais!”", M + 0.25, 1.6, cwm - 0.5, 1.0, size=15, color=CINZA, bold=True, line=1.25)
card(s, M + cwm + 0.22, 1.0, cwm, 1.75)
text(s, "COM A MARCA LIDA ANTES", M + cwm + 0.47, 1.2, cwm - 0.5, 0.25, size=9, color=AMARELO, font=MONO, space=1.2)
text(s, "“Toda taxa zero tem prazo de validade. Aqui a taxa fica na mesa — e quem faz a conta somos nós.”",
     M + cwm + 0.47, 1.6, cwm - 0.5, 1.1, size=14, bold=True, line=1.22)
text(s, [[("Ele não chutou o ângulo: leu que a dor é hora perdida na planilha e taxa que só aparece depois. ", DIM, False),
          ("E não para no vídeo — passa a ler o funil e agir onde ele vaza.", AMARELO, False)]],
     M, 3.0, W - 2 * M, 0.6, size=12, color=DIM, line=1.3)
FUN = [("01", "Interação", "Ninguém para no vídeo?", "Reescreve o gancho."),
       ("02", "Clique", "Assiste e não clica?", "Troca a promessa."),
       ("03", "Landing page", "Entra e sai na hora?", "Testa outra página."),
       ("04", "Oferta", "Chega e não fecha?", "Reposiciona a oferta."),
       ("05", "Venda", "A tag devolve o sinal.", "O diagnóstico aprende.")]
cwf = (W - 2 * M - 4 * 0.16) / 5
for i, (n, g, q, a) in enumerate(FUN):
    x = M + i * (cwf + 0.16)
    ult = i == len(FUN) - 1
    card(s, x, 3.85, cwf, 2.3, fill=AMARELO if ult else CARD)
    fg = PRETO if ult else CRU
    text(s, n, x + 0.2, 4.05, 1, 0.2, size=8, color=PRETO if ult else CINZA, font=MONO, space=1.2)
    text(s, g, x + 0.2, 4.32, cwf - 0.4, 0.35, size=14, bold=True, color=fg)
    text(s, q, x + 0.2, 4.85, cwf - 0.4, 0.6, size=10, color=PRETO if ult else DIM, line=1.25)
    text(s, a, x + 0.2, 5.5, cwf - 0.4, 0.55, size=10.5, bold=True, color=PRETO if ult else AMARELO, line=1.25)
s.notes_slide.notes_text_frame.text = "1:40–2:20 — O funil mostra a AÇÃO do agente, não só o nome da etapa."

# --------------------------------------------------------- 6 · negócio
s = slide()
label(s, "Quem paga — e por que não copiam", M, 0.6, 8)
text(s, "O contrato é com a rede, não com a loja.", M, 1.05, 10.5, 0.8, size=30, bold=True)
ROWS = [("Cliente", "Rede e franquia de varejo local. Vende uma vez, atende N unidades."),
        ("Tamanho", "Só o iFood tem 500 mil estabelecimentos em 1.500 cidades. É o piso."),
        ("Alternativa hoje", "Quatro fornecedores — pesquisa, redator, produtora, tráfego. Ou uma planilha vendida no YouTube."),
        ("Modelo", "Assinatura por unidade/mês."),
        ("Defensibilidade", "Cada campanha devolve o que converteu, para quem e com qual mecanismo. Copia-se o prompt numa tarde; não se copia o histórico.")]
y = 2.35
for k, v in ROWS:
    text(s, k.upper(), M, y, 2.6, 0.35, size=9.5, color=CINZA, font=MONO, space=1.2)
    text(s, v, M + 2.8, y - 0.04, W - 2 * M - 2.8, 0.8, size=13, line=1.3)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(M), In(y + 0.62), In(W - 2 * M), In(0.012))
    ln.fill.solid(); ln.fill.fore_color.rgb = C(0x2A, 0x2A, 0x28)
    ln.line.fill.background(); ln.shadow.inherit = False
    y += 0.92
s.notes_slide.notes_text_frame.text = "2:20–2:45 — Definir o valor da assinatura antes do pitch."

# ------------------------------------------------------------- 7 · fecho
s = slide(bg=AMARELO)
o = s.shapes.add_shape(MSO_SHAPE.OVAL, In(M), In(1.5), In(0.55), In(0.55))
o.fill.solid(); o.fill.fore_color.rgb = PRETO
o.line.fill.background(); o.shadow.inherit = False
text(s, "brandloop", M + 0.72, 1.55, 4, 0.5, size=26, bold=True, color=PRETO)
text(s, "A IA que aprende a sua marca antes de falar por ela.",
     M, 2.6, 10.2, 2.2, size=40, bold=True, color=PRETO, line=1.05)
text(s, "Não somos a agência do iFood — ele é o caso mais difícil do Brasil agora, "
        "e vocês conseguem verificar tudo o que o agente concluiu.",
     M, 5.1, 9.2, 1.0, size=14, color=C(0x2C, 0x2C, 0x22), line=1.35)
s.notes_slide.notes_text_frame.text = "2:45–3:00 — Fecho. Benchmark, não cliente."

# ------------------------------------------------------------ 8 · backup
s = slide()
label(s, "Backup — não apresentar · abrir se perguntarem", M, 0.6, 9)
text(s, "O que é real e o que é simulado", M, 1.05, 10.5, 0.8, size=30, bold=True)
LINHAS = [("REAL", MENTA, "Contexto e entrevista — /api/ingest e /api/interview, adapter Anthropic + OpenAI"),
          ("REAL", MENTA, "Análise de mercado — /api/research, Gorilla API, 1.516 conversas, search_id verificável"),
          ("REAL", MENTA, "Roteiro e criativo — /api/copy e /api/video, providers plugáveis"),
          ("SIMULADO", CINZA, "Publicação e campanha — simulado: true forçado pelo contrato de tipos"),
          ("SIMULADO", CINZA, "Performance — ROAS de 8 semanas não cabe em 5 horas")]
y = 2.3
for tag, cor, txt in LINHAS:
    text(s, tag, M, y, 1.5, 0.3, size=9, color=cor, font=MONO, bold=True, space=1.2)
    text(s, txt, M + 1.7, y - 0.05, W - 2 * M - 1.7, 0.6, size=12.5, line=1.3)
    y += 0.75
text(s, [[("A performance roda em ambiente de simulação — o mesmo harness que usaríamos para evals em "
           "produção. ", DIM, False),
          ("A diferença entre esconder e escolher é o que estamos mostrando aqui.", AMARELO, False)]],
     M, 6.25, W - 2 * M, 0.8, size=13, color=DIM, line=1.35)
s.notes_slide.notes_text_frame.text = "Slide de backup. Só abrir se perguntarem o que é real."

prs.save(OUT)
print("gerado:", OUT, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
